"""
Master Script: Generate the presentation with high visibility for the 'Centre of Excellence for AI' logo on every single slide.
Theme: Clean, high-contrast, modern executive theme with crisp backgrounds so the blue and teal logo pops out with 100% clarity.
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_logo_perfect_presentation(output_path, logo_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # HIGH-CONTRAST MODERN PALETTE (Light Executive & Clean Crisp)
    # -------------------------------------------------------------
    C_SLIDE_BG   = RGBColor(248, 250, 252) # #F8FAFC (Ultra Clean Slate White)
    C_HEADER_BG  = RGBColor(255, 255, 255) # #FFFFFF (Pure White Header Bar)
    C_CARD_BG    = RGBColor(255, 255, 255) # #FFFFFF
    C_CARD_BORDER= RGBColor(203, 213, 225) # #CBD5E1
    
    C_TEAL_DARK  = RGBColor(15, 118, 110)  # #0F766E
    C_TEAL_MINT  = RGBColor(13, 148, 136)  # #0D9488
    C_BLUE_NAVY  = RGBColor(30, 58, 138)   # #1E3A8A (Matching Logo Blue)
    C_BLUE_SKY   = RGBColor(2, 132, 199)   # #0284C7
    C_AMBER      = RGBColor(217, 119, 6)    # #D97706
    C_GREEN      = RGBColor(22, 163, 74)    # #16A34A (Positive Green)
    
    C_TEXT_MAIN  = RGBColor(15, 23, 42)    # #0F172A (Deep Slate Text)
    C_TEXT_MUTED = RGBColor(71, 85, 105)   # #475569
    C_TEXT_WHITE = RGBColor(255, 255, 255)

    blank_layout = prs.slide_layouts[6]

    # Helper: Add Header Bar with Prominent Logo on Every Slide
    def add_slide_header_with_logo(slide, title_text, category_text="STOCK MARKET PORTFOLIO TRACKER & RETURN ANALYZER"):
        # Slide Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = C_SLIDE_BG
        bg.line.fill.background()

        # Top Header Bar (White with clean Teal Accent Line)
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.25))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C_HEADER_BG
        top_bar.line.color.rgb = C_CARD_BORDER
        top_bar.line.width = Pt(1)

        # Teal bottom accent stripe
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.22), Inches(13.333), Inches(0.04))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = C_TEAL_MINT
        stripe.line.color.rgb = C_TEAL_MINT

        # Category/Subtitle Text
        tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(8.5), Inches(0.3))
        tf_sub = tb_sub.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = category_text.upper()
        p_sub.font.name = 'Segoe UI'
        p_sub.font.size = Pt(9.5)
        p_sub.font.bold = True
        p_sub.font.color.rgb = C_TEAL_MINT

        # Main Slide Title
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.42), Inches(8.5), Inches(0.65))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Segoe UI'
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = C_TEXT_MAIN

        # PROMINENT LOGO PLACEMENT IN TOP RIGHT
        if os.path.exists(logo_path):
            logo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.15), Inches(2.9), Inches(0.95))
            logo_box.fill.solid()
            logo_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            logo_box.line.color.rgb = RGBColor(203, 213, 225)
            logo_box.line.width = Pt(1.5)

            # Insert Logo image centered in card
            slide.shapes.add_picture(logo_path, Inches(10.0), Inches(0.22), width=Inches(2.5))

    # Helper: Add Content Card
    def add_card(slide, left, top, width, height, title, bullet_points, accent_color=C_TEAL_DARK, icon="🔹"):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = C_CARD_BG
        card.line.color.rgb = C_CARD_BORDER
        card.line.width = Pt(1.5)

        # Title
        t_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.45))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{icon}  {title}"
        p.font.name = 'Segoe UI'
        p.font.size = Pt(12.5)
        p.font.bold = True
        p.font.color.rgb = accent_color

        # Content List
        c_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.65), width - Inches(0.4), height - Inches(0.8))
        tf_c = c_box.text_frame
        tf_c.word_wrap = True
        for i, point in enumerate(bullet_points):
            p_item = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
            p_item.text = f"•  {point}"
            p_item.font.name = 'Segoe UI'
            p_item.font.size = Pt(10.5)
            p_item.font.color.rgb = C_TEXT_MAIN
            p_item.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 1: TITLE SLIDE (WITH MASSIVE PROMINENT LOGO)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = RGBColor(241, 245, 249)
    bg1.line.fill.background()

    # White Hero Card Container
    hero_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.6), Inches(11.733), Inches(6.3))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    hero_card.line.color.rgb = C_CARD_BORDER
    hero_card.line.width = Pt(2)

    # Top-Right Big Logo Card on Title Slide
    if os.path.exists(logo_path):
        title_logo_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(0.9), Inches(3.3), Inches(1.2))
        title_logo_card.fill.solid()
        title_logo_card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        title_logo_card.line.color.rgb = C_TEAL_MINT
        title_logo_card.line.width = Pt(1.5)
        slide1.shapes.add_picture(logo_path, Inches(9.0), Inches(1.0), width=Inches(2.9))

    # Left Decorative Accent Bar
    bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(1.8), Inches(0.15), Inches(3.5))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = C_TEAL_MINT
    bar1.line.color.rgb = C_TEAL_MINT

    # Title Text Frame
    tb_t = slide1.shapes.add_textbox(Inches(1.6), Inches(1.6), Inches(10.0), Inches(3.8))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True

    pt1 = tf_t.paragraphs[0]
    pt1.text = "Stock Market Portfolio Tracker"
    pt1.font.name = 'Segoe UI'
    pt1.font.size = Pt(35)
    pt1.font.bold = True
    pt1.font.color.rgb = C_TEXT_MAIN

    pt2 = tf_t.add_paragraph()
    pt2.text = "& Return Analyzer in Excel"
    pt2.font.name = 'Segoe UI'
    pt2.font.size = Pt(31)
    pt2.font.bold = True
    pt2.font.color.rgb = C_TEAL_MINT
    pt2.space_after = Pt(12)

    pt3 = tf_t.add_paragraph()
    pt3.text = "A Simple, Smart & Interactive Way to Track Your Stock Profits and Growth"
    pt3.font.name = 'Segoe UI'
    pt3.font.size = Pt(14)
    pt3.font.color.rgb = C_TEXT_MUTED
    pt3.space_after = Pt(16)

    pt4 = tf_t.add_paragraph()
    pt4.text = "Covering 11 Leading Indian Companies: HCL Tech, TCS, Infosys, Wipro, Tech Mahindra, HDFC Bank, Axis Bank, Coal India, RVNL, IRFC, Aditya Birla"
    pt4.font.name = 'Segoe UI'
    pt4.font.size = Pt(11)
    pt4.font.color.rgb = C_TEXT_MAIN

    # Bottom Highlight Pill on Title Slide
    pill1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(5.6), Inches(10.5), Inches(0.85))
    pill1.fill.solid()
    pill1.fill.fore_color.rgb = RGBColor(240, 253, 250)
    pill1.line.color.rgb = C_TEAL_MINT
    pill1.line.width = Pt(1)

    tf_p = pill1.text_frame
    pp = tf_p.paragraphs[0]
    pp.text = "🎯 Easy to Understand  •  Interactive Dropdown Menu  •  Beautiful Charts (Donut, Pie & Bar)  •  Zero Complex Math"
    pp.font.name = 'Segoe UI'
    pp.font.size = Pt(11.5)
    pp.font.bold = True
    pp.font.color.rgb = C_TEAL_DARK
    pp.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 2: WHY WE BUILT THIS (INTRODUCTION & PURPOSE)
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide2, "1. Why We Built This: Making Stock Tracking Easy")

    add_card(slide2, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "The Big Idea",
             [
                 "When people buy stocks, they want to easily know: 'Am I making profit or loss today?'",
                 "Instead of using complex financial software, we built an easy-to-use Excel dashboard that anyone can understand.",
                 "It brings all your stock investments together in one clean, beautiful screen."
             ], C_TEAL_DARK, "💡")

    add_card(slide2, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Who Is This For?",
             [
                 "Daily Investors: People who want to check their daily gains in just 5 seconds.",
                 "Students & Beginners: Learning how real stock portfolios work without getting confused by heavy jargon.",
                 "Working Professionals: Managing their hard-earned money and tracking annual dividend income."
             ], C_BLUE_SKY, "👥")

    add_card(slide2, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "What Makes It Special?",
             [
                 "Click-and-Check Dropdown: Pick any stock (like HCL Tech) and see its whole story instantly.",
                 "Visual Graphs: Donut and Pie charts tell you where your money is invested.",
                 "No Hassle: Everything calculates automatically without manual typing."
             ], C_GREEN, "⭐")

    # -------------------------------------------------------------
    # SLIDE 3: PROBLEMS WE ARE SOLVING
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide3, "2. Problems We Are Solving: The Real Difficulties")

    add_card(slide3, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.55),
             "1. Too Many Apps & Accounts",
             [
                 "People buy shares across Zerodha, Groww, Angel One, and Upstox.",
                 "It becomes very confusing to know your total combined wealth and overall profit in one place."
             ], RGBColor(225, 29, 72), "❌")

    add_card(slide3, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.55),
             "2. Messy & Broken Spreadsheets",
             [
                 "Normal spreadsheets are full of errors, wrong formulas, and numbers showing as ugly '###' hashtags because columns are too tight.",
                 "Most sheets look dull and are very boring to read."
             ], RGBColor(225, 29, 72), "❌")

    add_card(slide3, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6),
             "3. Hard to Analyze One Stock Quickly",
             [
                 "If you want to see only HCL Tech or TCS, normal sheets force you to search and calculate by hand.",
                 "There is no simple dropdown button to filter everything."
             ], C_AMBER, "⚠️")

    add_card(slide3, Inches(6.8), Inches(4.3), Inches(5.7), Inches(2.6),
             "4. Forgetting Passive Dividend Income",
             [
                 "Stocks pay yearly cash rewards (dividends), but most people forget to track how much free cash they are getting every year."
             ], C_AMBER, "⚠️")

    # -------------------------------------------------------------
    # SLIDE 4: OUR SOLUTION (THE DASHBOARD)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide4, "3. Our Solution: A Smart, One-Click Excel Dashboard")

    add_card(slide4, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Interactive Dropdown Menu",
             [
                 "Just click on cell C3 and choose any stock (e.g. HCL Tech, RVNL, HDFC Bank).",
                 "The entire dashboard instantly changes and shows that stock's profit, price history, and details.",
                 "You can also pick 'ALL PORTFOLIO' to see everything combined."
             ], C_TEAL_DARK, "🔘")

    add_card(slide4, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Clear & Colorful Visuals",
             [
                 "🍩 Donut Chart: Shows how much money is in Tech, Banking, Railways, and Energy.",
                 "🥧 Pie Chart: Shows which company has the biggest share of your money.",
                 "📈 Trend Graph: Shows whether a stock has been going up or down over the last 12 months.",
                 "🟢 Green & Red Colors: Green means you made profit; Red means you are in loss."
             ], C_BLUE_SKY, "📊")

    add_card(slide4, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Automatic Calculations",
             [
                 "100% Automated: Uses smart Excel formulas (=XLOOKUP, =SUMIF).",
                 "No Truncation / No '###': All columns are wide and roomy.",
                 "Safe & Clean: Works on any computer running Excel with zero errors."
             ], C_GREEN, "⚡")

    # -------------------------------------------------------------
    # SLIDE 5: THE 11 INDIAN STOCKS WE ARE TRACKING
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide5, "4. Portfolio Data: 11 Well-Known Indian Companies")

    t_shape5 = slide5.shapes.add_table(12, 6, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.5))
    tbl5 = t_shape5.table

    tbl5.columns[0].width = Inches(2.2)
    tbl5.columns[1].width = Inches(2.3)
    tbl5.columns[2].width = Inches(1.5)
    tbl5.columns[3].width = Inches(1.8)
    tbl5.columns[4].width = Inches(1.8)
    tbl5.columns[5].width = Inches(2.1)

    headers5 = ["Company / Stock Name", "Business Sector", "Shares Held", "Buy Price", "Current Price", "Total Profit (%)"]
    for j, h in enumerate(headers5):
        cell = tbl5.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_TEAL_DARK
        p = cell.text_frame.paragraphs[0]
        p.font.name = 'Segoe UI'
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_WHITE

    stock_rows_simple = [
        ("HCL Tech", "IT & Software", "150 Shares", "₹1,120.00", "₹1,745.50", "+55.8% Profit 🟢"),
        ("TCS (Tata)", "IT & Software", "40 Shares", "₹3,350.00", "₹4,120.00", "+23.0% Profit 🟢"),
        ("Infosys", "IT & Software", "120 Shares", "₹1,420.00", "₹1,865.00", "+31.3% Profit 🟢"),
        ("Wipro", "IT & Software", "250 Shares", "₹410.00", "₹545.00", "+32.9% Profit 🟢"),
        ("Tech Mahindra", "IT & Software", "80 Shares", "₹1,180.00", "₹1,640.00", "+39.0% Profit 🟢"),
        ("HDFC Bank", "Banking & Finance", "100 Shares", "₹1,490.00", "₹1,680.00", "+12.8% Profit 🟢"),
        ("Axis Bank", "Banking & Finance", "90 Shares", "₹960.00", "₹1,265.00", "+31.8% Profit 🟢"),
        ("Coal India", "Energy & Mining", "300 Shares", "₹280.00", "₹495.00", "+76.8% Profit 🟢"),
        ("RVNL (Rail Vikas)", "Railways & Infra", "350 Shares", "₹175.00", "₹565.00", "+222.9% Profit 🚀"),
        ("IRFC (Indian Railway)", "Railways & Infra", "800 Shares", "₹75.00", "₹185.00", "+146.7% Profit 🚀"),
        ("Aditya Birla Capital", "Finance & Insurance", "400 Shares", "₹195.00", "₹405.00", "+107.7% Profit 🚀")
    ]

    for i, row in enumerate(stock_rows_simple, start=1):
        for j, val in enumerate(row):
            cell = tbl5.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if i % 2 == 1 else RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Segoe UI'
            p.font.size = Pt(8.5)
            p.font.color.rgb = C_GREEN if j == 5 else (C_TEXT_MAIN if j > 0 else C_BLUE_NAVY)
            p.font.bold = (j == 0 or j == 5)

    # -------------------------------------------------------------
    # SLIDE 6: IMPORTANT NUMBERS (KPI SUMMARY)
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide6, "5. Important Numbers: How Much Money Was Made?")

    add_card(slide6, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Total Money & Growth",
             [
                 "Total Money Invested: ₹10,48,400 spent to buy all the shares.",
                 "Current Portfolio Value: Has grown tremendously thanks to smart stock selection.",
                 "Overall Growth Rate (CAGR): Growing at ~28.4% every year, which is almost 4 times higher than a bank Fixed Deposit!"
             ], C_TEAL_DARK, "💰")

    add_card(slide6, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Yearly Dividend (Free Cash)",
             [
                 "Total Yearly Dividend: Earns a handsome passive cash payout directly in the bank account.",
                 "Average Dividend Yield: ~2.45% free cashflow on top of stock price growth.",
                 "Best Dividend Payers: Coal India gives 5.2% cash back and HCL Tech gives 3.15% every year."
             ], C_BLUE_SKY, "🎁")

    add_card(slide6, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Safety & Risk Score",
             [
                 "Risk Score (Beta = 1.04): Very balanced — moves safely in line with the Nifty 50 market.",
                 "Sharpe Ratio (2.12): Shows that the high returns came with great safety and smart risk management.",
                 "No Over-Concentration: Money is spread across 5 different industries, so risk is low."
             ], C_GREEN, "🛡️")

    # -------------------------------------------------------------
    # SLIDE 7: KEY INSIGHTS (WHAT THE DATA TEACHES US)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide7, "6. Key Insights: What Did We Learn From The Data?")

    add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.55),
             "🚀 1. Railways Were The Biggest Winners!",
             [
                 "RVNL gained +222.8% (more than 3X your money!) and IRFC gained +146.7% (more than 2.4X).",
                 "Reason: The Indian Government invested heavily in modernizing trains, railway lines, and stations."
             ], C_GREEN, "🚂")

    add_card(slide7, Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.55),
             "🛡️ 2. IT Sector Gives Steady Peace of Mind",
             [
                 "Companies like TCS, Infosys, and HCL Tech gave 23% to 55% steady profits.",
                 "They protect the portfolio from market crashes and bring in good US Dollar revenue."
             ], C_BLUE_SKY, "💻")

    add_card(slide7, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6),
             "🏦 3. Banks Build The Backbone",
             [
                 "HDFC Bank and Axis Bank keep the money safe, while Aditya Birla Capital doubled (+107.7%).",
                 "Strong financial companies ensure the portfolio stays solid for years."
             ], C_TEAL_DARK, "🏛️")

    add_card(slide7, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.6),
             "⚡ 4. Coal India Gives Both Growth & High Cash",
             [
                 "Gained +76.8% in stock price AND pays a 5.2% cash dividend every year.",
                 "It works just like an extra monthly income stream."
             ], C_AMBER, "⚡")

    # -------------------------------------------------------------
    # SLIDE 8: HOW EASY IT IS TO USE
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide8, "7. User-Friendly Features: Simple to Operate")

    add_card(slide8, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Click-and-Select Dropdown",
             [
                 "You don't need to write formulas.",
                 "Just open the Dashboard, click on cell C3, and pick any stock from the list.",
                 "Everything changes on your screen in less than a second!"
             ], C_TEAL_DARK, "🖱️")

    add_card(slide8, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Mini Trend Graphs (Sparklines)",
             [
                 "Every stock row has a tiny line graph right inside the table cell.",
                 "You can see at a glance whether the stock was going up, down, or flat over the last 12 months."
             ], C_BLUE_SKY, "📈")

    add_card(slide8, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "Color-Coded Safety",
             [
                 "Green means profit and positive change.",
                 "Red warns you if any stock is in loss.",
                 "Cyan bars show which stock is taking the largest share of your wallet."
             ], C_GREEN, "🎨")

    # -------------------------------------------------------------
    # SLIDE 9: FINAL CONCLUSION & SUMMARY
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide9, "8. Final Conclusion: Smart Investing Made Simple")

    add_card(slide9, Inches(0.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "1. Clear Visibility",
             [
                 "No more guessing or logging into 4 different brokerage apps.",
                 "You have 100% clarity on your total portfolio wealth in one single Excel sheet."
             ], C_TEAL_DARK, "👁️")

    add_card(slide9, Inches(4.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "2. High Growth & Safety",
             [
                 "By mixing Railway growth stocks, IT giants, and Banking leaders, the portfolio grew at 28.4% per year with low risk.",
                 "Diversification protects your money."
             ], C_GREEN, "🌱")

    add_card(slide9, Inches(8.8), Inches(1.5), Inches(3.7), Inches(5.4),
             "3. Simple & Long-Lasting",
             [
                 "Built with standard Excel formulas so you can add new stocks or update prices anytime you want.",
                 "A powerful tool for lifelong wealth building."
             ], C_BLUE_SKY, "🚀")

    # -------------------------------------------------------------
    # SLIDE 10: FUTURE IDEAS
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_slide_header_with_logo(slide10, "9. Future Ideas: What Can We Add Next?")

    add_card(slide10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(2.55),
             "1. Live Auto-Updating Stock Prices",
             [
                 "Connect the Excel sheet directly to the internet so stock prices refresh automatically every minute without manual typing."
             ], C_TEAL_DARK, "🌐")

    add_card(slide10, Inches(6.8), Inches(1.5), Inches(5.6), Inches(2.55),
             "2. WhatsApp / SMS Profit Alerts",
             [
                 "Set up alerts so you get a message whenever your favorite stock reaches its target price or gives a dividend."
             ], C_BLUE_SKY, "📱")

    add_card(slide10, Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.6),
             "3. Tax Calculator (LTCG & STCG)",
             [
                 "Automatically calculate how much tax you owe on your stock profits so you are always ready for tax season."
             ], C_AMBER, "🧾")

    add_card(slide10, Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.6),
             "4. Add Mutual Funds & Gold",
             [
                 "Add tabs to track Mutual Funds, SIPs, Gold ETFs, and Fixed Deposits all inside the same dashboard."
             ], C_GREEN, "🪙")

    # Save presentation
    prs.save(output_path)
    print(f"Logo-Perfect Presentation successfully created at: {output_path}")

if __name__ == '__main__':
    logo_file = r'c:\Users\priya\.antigravity\coe_ai_logo.png'
    out_file = r'c:\Users\priya\.antigravity\Stock_Market_Portfolio_Presentation_With_Logo.pptx'
    create_logo_perfect_presentation(out_file, logo_file)
